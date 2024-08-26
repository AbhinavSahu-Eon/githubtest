

import os

packages = [
    "PyPDF2",
    "langchain_community",
    "tiktoken",
    "langchain-openai",
    "langchainhub",
    "chromadb",
    "langchain",
    "openai",
    "PyPDF",
    "gradio",
   # "neo4j",
    "py2neo",
    "python-docx",
    "python-pptx"
]

for package in packages:
    os.system(f"python -m pip install {package}")

import gradio as gr
import os
from dotenv import load_dotenv
#from py2neo import Graph
#from langchain_community.vectorstores import Neo4jVector
#from langchain_community.graphs import Neo4jGraph
from langchain.text_splitter import CharacterTextSplitter
from langchain.document_loaders import PyPDFLoader
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain import OpenAI
from langchain.embeddings import OpenAIEmbeddings
from langchain_openai import ChatOpenAI
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnablePassthrough
from langchain_community.document_loaders import WebBaseLoader
from langchain import hub
from langchain_community.document_loaders import Docx2txtLoader
from langchain_community.document_loaders import TextLoader
from pptx import Presentation
import PyPDF2

load_dotenv()

# Define functions for loading API keys and PDF processing
def load_api_key(file_path):
    try:
        with open(file_path, 'r') as file:
            api_key = file.read().strip()
            return api_key
    except FileNotFoundError:
        print("Error: File not found.")
    except Exception as e:
        print(f"Error: {e}")


def create_vector_index(driver, dimension: int) -> None:
    index_query = "CALL db.index.vector.createNodeIndex('stackoverflow', 'Question', 'embedding', $dimension, 'cosine')"
    try:
        driver.query(index_query, {"dimension": dimension})
    except:  # Already exists
        pass
    index_query = "CALL db.index.vector.createNodeIndex('top_answers', 'Answer', 'embedding', $dimension, 'cosine')"
    try:
        driver.query(index_query, {"dimension": dimension})
    except:  # Already exists
        pass


# Load API keys

file_path_openai = "C:/Users/kesha/Downloads/GenAI_Priyanka/Data/api_key.txt"
file_path_langchain = "C:/Users/kesha/Downloads/GenAI_Priyanka/Data/langchain-api.txt"
openai_api_key = load_api_key(file_path_openai)   
langchain_api_key = load_api_key(file_path_langchain)  

os.environ['OPENAI_API_KEY'] = openai_api_key
os.environ['LANGCHAIN_API_KEY'] = langchain_api_key

# Define the directory containing PDF files

#directory = "C:/Users/kesha/Downloads/GenAI_Priyanka/Data/PDFs"

# Initialize the vector store
from langchain.schema import Document
import docx

vectorstore = Chroma()

############ TEXT LOADERS ############
# Functions to read different file types

class DocumentLoader: # For .txt files
    def __init__(self, file_path):
        self.file_path = file_path

    def load(self):
        doc = docx.Document(self.file_path)
        line_number = 0
        page_number = 1  # Page number starts from 1

        for paragraph in doc.paragraphs:

            yield Document(
                page_content=paragraph.text,
                metadata={"line_number": line_number, "page_number": page_number, "source": self.file_path},
            )
            line_number += 1

class PptxLoader: # For .ppt or .pptx files
    def __init__(self, file_path):
        self.file_path = file_path

    def load(self):
        prs = Presentation(self.file_path)
        slide_number = 0

        for slide in prs.slides:
            slide_number += 1
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    yield Document(
                        page_content=shape.text,
                        metadata={"slide_number": slide_number, "source": self.file_path},
                    )

def read_pdf(file_path):
    loader = PyPDFLoader(file_path)
    text = loader.load()
    return text

def read_word(file_path):
    loader = DocumentLoader(file_path)
    text = loader.load()
    return text

def read_txt(file_path):    
    loader = TextLoader(file_path)
    text = loader.load()
    return text

def read_ppt(file_path):
    loader = PptxLoader(file_path)
    text = loader.load()
    return text

train_directory = 'C:/Users/kesha/Downloads/GenAI_Priyanka/Data/PDFs'

for filename in os.listdir(train_directory):
    file_path = os.path.join(train_directory, filename)
    if filename.endswith(".pdf"):
        combined_text = read_pdf(file_path)
        # Split the PDF text into chunks
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=50)
        all_splits = text_splitter.split_documents(combined_text)

        # Wrap split texts in Document objects
        #documents = [Document(page_content=chunk) for chunk in all_splits]
        
        vectorstore = Chroma.from_documents(documents=all_splits,
                                    embedding=OpenAIEmbeddings())

    elif filename.endswith(".docx") or filename.endswith(".doc"):
        combined_text = read_word(file_path)
        # Split the PDF text into chunks
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=50)
        all_splits = text_splitter.split_documents(combined_text)

        # Wrap split texts in Document objects
        #documents = [Document(page_content=chunk) for chunk in all_splits]
        
        vectorstore = Chroma.from_documents(documents=all_splits,
                                    embedding=OpenAIEmbeddings())
    elif filename.endswith(".txt"):
        combined_text = read_txt(file_path)
        # Split the PDF text into chunks
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=50)
        all_splits = text_splitter.split_documents(combined_text)

        # Wrap split texts in Document objects
        #documents = [Document(page_content=chunk) for chunk in all_splits]
        
        vectorstore = Chroma.from_documents(documents=all_splits,
                                    embedding=OpenAIEmbeddings())
    
    elif filename.endswith(".ppt") or filename.endswith(".pptx"):
        combined_text = read_ppt(file_path)
        # Split the PDF text into chunks
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=50)
        all_splits = text_splitter.split_documents(combined_text)

        # Wrap split texts in Document objects
        #documents = [Document(page_content=chunk) for chunk in all_splits]
        
        vectorstore = Chroma.from_documents(documents=all_splits,
                                    embedding=OpenAIEmbeddings())
##################################################################################################





#vectorstore = Chroma.from_documents(documents=splits, embedding=OpenAIEmbeddings())


# Define the RAG framework components
retriever = vectorstore.as_retriever()
prompt = hub.pull("rlm/rag-prompt")
llm = ChatOpenAI(model_name="gpt-3.5-turbo", temperature=0)

# Post-processing function
def format_docs(docs):
    return "\n\n".join(doc.page_content for doc in docs)

# Define the RAG chain
rag_chain = (
    {"context": retriever | format_docs, "question": RunnablePassthrough()}
    | prompt
    | llm
    | StrOutputParser()
)


# Function to generate response
def generate_response(Question):
    response = rag_chain.invoke(Question)
    return response

# Define your function to generate responses using your machine learning model
def generate_response(Question):
    # Define basic greetings and responses
    greetings = ["hi!", "hello!", "hey!", "Hola!", "hi", "hello", "hey", "Hola"]
    thank_you = ["thank you", "thanks", "thanks!"]

    # Convert the input question to lowercase for case-insensitive matching
    question = Question.lower()

    # Check if the input question is a basic greeting
    if question in greetings:
        response = "Hello! How can I assist you today?"
    # Check if the input question is an expression of gratitude
    elif question in thank_you:
        response = "You're welcome! If you have any more questions or if there's anything else I can assist you with, please don't hesitate to let me know."
    else:
        # Call your machine learning model here to generate a response based on the query
        response = rag_chain.invoke(question)
    return response


# Create the Gradio interface with custom CSS
iface = gr.Interface(
    fn=generate_response,
    inputs="text",
    outputs="text",
    title="<div style='padding-left: 35px; font-family: Zefani; font-size: 40px;'>NewYork Life</div>",
    description="<div style='padding-left: 470px; font-size: x-large; font-family: Zefani; padding-left: 352px; '>Your personalized insurance advisor</div> <br><img src='https://chambermaster.blob.core.windows.net/images/customers/2402/members/12917/logos/MEMBER_PAGE_HEADER/nyl_bgal_logo_color_rgb.png' style='width: 850px; padding-left: 365px; height: 170px;'>",
    examples=[
        ["What happens if a player lands on an unowned property in Monopoly?"],
        ["How does a player get out of Jail in Monopoly?"],
        ["How do players score points in Ticket to Ride?"]
    ],
    #theme="compact",    # Use the compact theme for buttons
    allow_flagging="manual",  # Allow flagging examples
)

# Launch the Gradio interface  
# GIVE A NEW PORT NUMBER EVERYTIME THIS PROGRAM IS RUN
iface.launch(server_port=1075, share=True)                     